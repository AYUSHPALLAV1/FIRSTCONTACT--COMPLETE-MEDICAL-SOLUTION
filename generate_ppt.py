from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- Helper Functions ---
    def set_font(paragraph, font_name='Arial', font_size=18, bold=False, color=None):
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color

    def add_slide(layout_index, title_text, content_text=None):
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        
        # Content
        if content_text and len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            body.text = content_text
        
        return slide

    # --- 1. Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title Layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "FirstContact"
    subtitle.text = "AI-Powered Healthcare System\n\nDigital Healthcare Accessibility and AI-Driven Preliminary Diagnosis"

    # --- 2. Agenda ---
    content = (
        "1. Project Overview\n"
        "2. Technical Approach & Algorithm\n"
        "3. Core Functionalities\n"
        "4. Literature Survey\n"
        "5. Conclusion"
    )
    add_slide(1, "Agenda", content)

    # --- 3. Project Overview ---
    content = (
        "Main Topic: Digital Healthcare Accessibility\n\n"
        "• Comprehensive healthcare platform\n"
        "• Bridges gap between patients and immediate medical assistance\n"
        "• Filters non-emergency cases to reduce hospital burden\n"
        "• Guides patients to appropriate specialists"
    )
    add_slide(1, "Project Overview", content)

    # --- 4. Technical Approach: Algorithm ---
    content = (
        "Algorithm: Random Forest Classifier (Ensemble Learning)\n\n"
        "• Library: sklearn.ensemble.RandomForestClassifier\n"
        "• Why Random Forest?\n"
        "  - High accuracy for tabular data\n"
        "  - Handles high-dimensional data (many symptoms)\n"
        "  - Resistant to overfitting compared to single decision trees\n"
        "• Training Data: 50,000 synthetic samples based on clinical reference matrix"
    )
    add_slide(1, "Technical Approach: Algorithm", content)

    # --- 5. Technical Approach: Stack ---
    content = (
        "Feature Engineering:\n"
        "• Numerical: Age, BMI, BP, Heart Rate, Blood Panel\n"
        "• Categorical: One-Hot Encoded Symptoms, History, Lifestyle\n\n"
        "Tech Stack:\n"
        "• Backend: Node.js (Express) + Python (Inference)\n"
        "• Frontend: HTML/CSS/JS (Responsive)\n"
        "• Database: MongoDB (User Profiles, Doctors, Logs)"
    )
    add_slide(1, "Technical Approach: Stack & Data", content)

    # --- 6. Core Functionalities: Medidose ---
    content = (
        "Module 1: Medidose (AI Diagnosis)\n\n"
        "• Input: Demographics, Vitals, Symptoms\n"
        "• Process: AI Analysis (Top 3 predictions)\n"
        "• Output: Preliminary diagnosis + Confidence Score\n"
        "• Goal: Triage (Home Care vs. Urgent Care)"
    )
    add_slide(1, "Use Case: Medidose", content)

    # --- 7. Core Functionalities: Consultation ---
    content = (
        "Module 2: Consultation (Telemedicine)\n\n"
        "• Doctor Booking: Browse profiles, check availability, book video calls\n"
        "• Payments: Integrated with Stripe\n"
        "• Hospital Locator: Google Maps API integration\n"
        "• Connection: Direct link from Diagnosis results to Booking"
    )
    add_slide(1, "Use Case: Consultation", content)

    # --- 8. Core Functionalities: Insurance ---
    content = (
        "Module 3: Insurance Policy Analysis\n\n"
        "• Feature: Upload Policy Document (PDF)\n"
        "• Analysis: Text processing to check disease coverage\n"
        "• Benefit: Financial transparency for patients\n"
        "• Logic: Matches diagnosed condition against policy inclusions/exclusions"
    )
    add_slide(1, "Use Case: Insurance", content)

    # --- 9. Literature Survey ---
    content = (
        "Foundational Research:\n\n"
        "• Breiman, L. (2001): Random Forests (Algorithm basis)\n"
        "• Jiang, F., et al. (2017): AI in Healthcare (Triage validity)\n"
        "• Topol, E. J. (2019): Human-AI Convergence (Clinical workflow)\n"
        "• Hwang, J. I., & Park, H. A. (2019): Symptom Checker Accuracy\n"
        "• Wootton, R. (2001): Telemedicine (Remote care validity)"
    )
    add_slide(1, "Literature Survey", content)

    # --- 10. Conclusion ---
    content = (
        "Summary:\n"
        "• FirstContact provides an end-to-end healthcare solution.\n"
        "• Combines AI prediction with practical medical access.\n\n"
        "Future Scope:\n"
        "• Deep Learning for Image Diagnosis (X-rays/Skin)\n"
        "• Real-time wearable integration\n"
        "• Multi-language support"
    )
    add_slide(1, "Conclusion & Future Scope", content)

    # Save
    output_path = "FirstContact_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
